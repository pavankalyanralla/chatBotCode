from flask_restful import Resource, reqparse
from database.queries import DB_Queries
from database.dB_Execute import *
from public.strings.strings import *
from database.queries import *
# from flask_mysqldb import MySQL
from src.server.constants.response import *
from datetime import datetime
import requests
from datetime import datetime, timedelta
import os.path
from googleapiclient.discovery import build
from google_auth_oauthlib.flow import InstalledAppFlow
from google.auth.transport.requests import Request
from google.oauth2.credentials import Credentials
import smtplib
import sys
from email.mime.multipart import MIMEMultipart
from email.mime.text import MIMEText
from email.mime.application import MIMEApplication
from flask import Flask, jsonify, request
from docx.api import Document
from openpyxl import load_workbook
from pptx import Presentation
from io import BytesIO
import os
import tempfile
import pdfplumber
import openai
import webbrowser


class addUserFormAndSuggestion(Resource):
    def post(self):
        requestData = request.get_json()
        if (not request.get_json()
                or not 'name' in request.get_json() or requestData['name'] == ''
                or not 'email' in request.get_json() or requestData['email'] == ''
                or not 'country' in request.get_json() or requestData['country'] == ''
                or not 'mobileNumber' in request.get_json()
                or not 'applicationType' in request.get_json() or requestData['applicationType'] == ''
                or not 'reqDescription' in request.get_json() or requestData['reqDescription'] == ''):
            return emptyDataResponse()
        name = requestData['name']
        country = requestData['country']
        email = requestData['email']
        mobileNumber = requestData['mobileNumber']
        applicationType = requestData['applicationType']
        reqDescription=requestData['reqDescription']
        # try:
        queryForGPT="conclude me in one line whether the user is looking for a job or service from a software company? if the user says,"+reqDescription
        filterJobOrservice=ask_gpt(queryForGPT)
        print("chat......", filterJobOrservice)
        createdTs = datetime.now()
        if (filterJobOrservice['statusCode'] == 200):
            insertComapyDetails = InsertIntoDB(DB_Queries.InsertFormDetails,name,country,email,mobileNumber,reqDescription,str(filterJobOrservice),applicationType,createdTs)
            # addDataInMongo = gptForm.insert_one({"name" : name,"country":country,"email":email,"mobile_number":mobileNumber,"application_type":applicationType,"req_description":reqDescription,"created_ts":createdTs})
            response = {
                "filterJobOrservicewithchatgpt" : filterJobOrservice['data'],
                "success" : True
                }
        else:
            response = {
                "message" : filterJobOrservice['data'],
                "success" : True
            }
        return response
        # except Exception as e:
        #     response = exceptionResponseWithMessage(str(e))

#ask chatgpt for synopsis of requirements
def ask_gpt(synopsis):
    api_endpoint = 'https://api.openai.com/v1/chat/completions'
    api_key='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'
    headers = {
        'Authorization': f'Bearer {api_key}',
        'Content-Type': 'application/json'
    }
    # Prepare the data payload with the prompt
    data = {
        "model":"gpt-3.5-turbo",
        "messages":[{"role":"user","content":synopsis}]}
    # Make a POST request to ChatGPT
    response = requests.post(api_endpoint, headers=headers, json=data)
    # Get the generated completion
    if response.status_code == 200:
        dataf = {'data':response.json()['choices'][0]['message']['content'],'statusCode':response.status_code}
    else:
        dataf = {'data':response.text,'statusCode':response.status_code}
        # return("Error:", response.status_code, response.text)
    return dataf

class synopsisFromUploadedFile(Resource):
    def post(self):
        redirecrUrl = ''
        if 'file' not in request.files:
            return 'Please select file and upload.'
        uploadedFile = request.files['file']
        if uploadedFile.filename == '':
            return 'File not selected.'
        else:
            temp_directory = tempfile.mkdtemp()
            temp_file_path = os.path.join(temp_directory, uploadedFile.filename)
            uploadedFile.save(temp_file_path)
            file_extension = os.path.splitext(uploadedFile.filename)[-1].lower()
            if file_extension == ".pdf" or file_extension == ".docx" or file_extension == ".xlsx" or file_extension == ".pptx":
                if file_extension == ".pdf":
                    with pdfplumber.open(temp_file_path) as pdf:
                        text = ''
                        for page in pdf.pages:
                            text += page.extract_text()
                elif file_extension == ".docx":
                    doc = Document(temp_file_path)
                    text = ""
                    for para in doc.paragraphs:
                        text += para.text + "\n"
                elif file_extension == ".xlsx":
                    text = ""
                    workbook = load_workbook(filename=temp_file_path)
                    for sheet_name in workbook.sheetnames:
                        sheet = workbook[sheet_name]
                        for row in sheet.iter_rows(values_only=True):
                            text += ' '.join(str(cell) for cell in row if cell) + '\n'
                elif file_extension == ".pptx":
                    presentation = Presentation(temp_file_path)
                    text = ""
                    for slide in presentation.slides:
                        for shape in slide.shapes:
                            if hasattr(shape, "text"):
                                text += shape.text + "\n"
            else:
                response = "Unsupported file format. Please upload 'pdf' or 'docx' or 'xlsx' or 'pptx' files."
            # Clean up: delete temporary file and directory
            os.remove(temp_file_path)
            os.rmdir(temp_directory)
            queryforGPT="give me synopsis in 100 to 250 words if the user says,"+text
            synopsisFromGPT=ask_gpt(queryforGPT)
            note="check the above synopsis and say in one word whether it is resume or service seeking or \"not related to resume and service seeking\""
            suggestionForRedirection=ask_gpt("synosis given as "+synopsisFromGPT+"\n"+note)
            # print("synosis given as "+synopsisFromGPT+"\n"+note)
            if suggestionForRedirection.lower()=="resume" or suggestionForRedirection.lower()=="job":
                redirecrUrl = "https://terralogic.com/careers/"
            elif suggestionForRedirection.lower()=="service seeking":
                redirecrUrl = "https://terralogic.com/contact-us/"
            return jsonify({"synopsisFromGPT": synopsisFromGPT,"suggestionForRedirection":suggestionForRedirection,
                            "redirectUrl":redirecrUrl})





